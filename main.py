"""

CLI tool to find factual/logical inconsistencies across slides in a PowerPoint (.pptx).
- Extracts text from slides
- Extracts images and runs OCR (pytesseract) on images
- Performs deterministic checks (numbers, percentages, dates/timelines, repeated metric mismatches)
- Optionally calls Gemini 2.5 Flash for semantic contradiction detection

Dependencies:
  pip install python-pptx pytesseract pillow regex python-dateutil google-genai tqdm

Notes:
  - You must have Tesseract installed for OCR (if using OCR). On Ubuntu: `sudo apt install tesseract-ocr`
  - Set GEMINI_API_KEY env var (or pass --api-key) to use Gemini.
  - If you can't/care not to use Gemini, run with --no-llm (only deterministic checks run).
"""

import os
import re
import argparse
import json
import tempfile
import shutil
from collections import defaultdict, Counter
from datetime import datetime
from dateutil import parser as dateparser
from pptx import Presentation
from PIL import Image
import pytesseract
from tqdm import tqdm

# Optional: Gemini client
USE_GEMINI = True
try:
    from google.genai import client as gemini_client
    from google.genai import types as gemini_types
except Exception:
    # We'll check later and tell user to install google-genai if they want LLM.
    gemini_client = None

# ---------------------------
# Utilities: extraction
# ---------------------------

NUM_RE = re.compile(r'(?<![\d.,])([+-]?\(?\d{1,3}(?:[,\d]{0,3})*(?:\.\d+)?\)?)')
PCT_RE = re.compile(r'([+-]?\d+(?:\.\d+)?)\s*%')
DATE_LIKE_RE = re.compile(r'\b(?:\d{1,2}[/\-]\d{1,2}[/\-]\d{2,4}|\d{4}|\b(?:Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)[a-z]*\s+\d{2,4})\b', re.IGNORECASE)

def normalize_number(s):
    # remove commas and parentheses
    s = s.strip()
    neg = False
    if s.startswith('(') and s.endswith(')'):
        neg = True
        s = s[1:-1]
    s = s.replace(',', '')
    try:
        if '.' in s:
            val = float(s)
        else:
            val = int(s)
        return -val if neg else val
    except:
        try:
            return float(s)
        except:
            return None

def extract_numbers_and_percents(text):
    nums = []
    for m in NUM_RE.finditer(text):
        n = normalize_number(m.group(1))
        if n is not None:
            nums.append((m.group(1), n))
    pcts = []
    for m in PCT_RE.finditer(text):
        try:
            p = float(m.group(1))
            pcts.append((m.group(0), p))
        except:
            pass
    return nums, pcts

def extract_dates(text):
    dates = []
    for m in DATE_LIKE_RE.finditer(text):
        snippet = m.group(0)
        try:
            dt = dateparser.parse(snippet, fuzzy=True, default=datetime(1900,1,1))
            dates.append((snippet, dt.date().isoformat()))
        except:
            pass
    return dates



def pptx_to_slides(pptx_path, images_outdir=None, ocr_on_images=True):
    prs = Presentation(pptx_path)
    slides_data = {}
    tmpdir = None
    if images_outdir is None:
        tmpdir = tempfile.mkdtemp(prefix="pptx_imgs_")
        images_outdir = tmpdir
    slide_no = 0
    for slide in prs.slides:
        slide_no += 1
        texts = []
        # shapes text
        for shape in slide.shapes:
            try:
                if shape.has_text_frame:
                    texts.append(shape.text)
                elif shape.shape_type.name == 'PICTURE':
                    # will extract below
                    pass
                elif shape.has_table:
                    # extract table text
                    try:
                        rows = []
                        for r in shape.table.rows:
                            row_texts = [c.text for c in r.cells]
                            rows.append(' | '.join(row_texts))
                        texts.append('\n'.join(rows))
                    except Exception:
                        pass
            except Exception:
                # some shapes may raise has_text_frame incorrectly
                pass

        # slide notes
        notes_text = ''
        try:
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text
        except Exception:
            notes_text = ''

        # images extraction
        imgs = []
        ocr_texts = []
        for i,shape in enumerate(slide.shapes):
            try:
                if shape.shape_type.name == 'PICTURE':
                    img = shape.image
                    ext = img.ext
                    img_bytes = img.blob
                    fname = os.path.join(images_outdir, f"slide{slide_no}_img{i}.{ext}")
                    with open(fname, "wb") as f:
                        f.write(img_bytes)
                    imgs.append(fname)
                    if ocr_on_images:
                        try:
                            pil = Image.open(fname)
                            ocr_t = pytesseract.image_to_string(pil)
                            if ocr_t.strip():
                                ocr_texts.append(ocr_t)
                        except Exception:
                            pass
            except Exception:
                pass

        bigtext = "\n\n".join([t for t in texts if t and t.strip()])
        bigocr = "\n\n".join([t for t in ocr_texts if t and t.strip()])
        slides_data[slide_no] = {
            "text": bigtext,
            "images": imgs,
            "ocr_text": bigocr,
            "notes": notes_text
        }
    return slides_data, tmpdir


def find_conflicting_numbers(slides_data, fuzzy_name_window=40):
    """
    Find repeated metric mentions with different numeric values across slides.
    Strategy:
      - For each slide, find "phrases" near numbers (e.g.,  "Revenue: $12,000" or "2024 revenue 12,000").
      - Build mapping phrase -> set(values, slide_no)
      - Flag phrase groups that have >1 distinct numeric values
    """
    phrase_values = defaultdict(list)  # phrase -> list of (value, slide_no, raw)
    for slide_no, s in slides_data.items():
        text = s["text"] + "\n" + s["ocr_text"] + "\n" + s["notes"]
        # find numbers and nearby words
        for m in NUM_RE.finditer(text):
            raw = m.group(1)
            val = normalize_number(raw)
            if val is None: continue
            start, end = m.span()
            # get phrase window
            left = text[max(0, start - fuzzy_name_window):start].strip()
            right = text[end:end + fuzzy_name_window].strip()
            # phrase candidate: last 40 chars left + first 40 right
            phrase = (left + " " + right).replace("\n", " ").strip()
            # also try to extract a short metric name from left using simple heuristics
            metric_name = None
            # look for word before colon or keyword before number
            colon_search = re.search(r'([A-Za-z ]{1,60}):\s*$', left)
            if colon_search:
                metric_name = colon_search.group(1).strip()
            else:
                # take up to last 5 words in left
                metric_name = " ".join(left.split()[-5:])
            phrase_values[metric_name].append((val, slide_no, raw, phrase))
    # evaluate contradictions
    contradictions = []
    for metric, vals in phrase_values.items():
        distinct_vals = set(v[0] for v in vals)
        if len(distinct_vals) > 1:
            contradictions.append({
                "metric_hint": metric,
                "values": [{"value": v[0], "slide": v[1], "raw": v[2], "context": v[3]} for v in vals]
            })
    return contradictions

def percentages_sum_checks(slides_data):
    """
    For items that look like parts/percent lists, check if they add to ~100.
    Strategy:
      - Find contiguous percentages in a slide (or in OCR text) that look like a distribution.
      - If there are >=2 percentages and sum deviates from 100 by a tolerance, flag it.
    """
    issues = []
    tol = 3.0  # percent tolerance
    for slide_no, s in slides_data.items():
        text = s["text"] + "\n" + s["ocr_text"]
        pcts = PCT_RE.findall(text)
        if len(pcts) >= 2:
            try:
                pvals = [float(p) for p in pcts]
                total = sum(pvals)
                if abs(total - 100.0) > tol:
                    issues.append({
                        "slide": slide_no,
                        "percentages": pvals,
                        "total": total,
                        "message": f"Percentages on slide {slide_no} sum to {total:.2f} (tolerance {tol}%)"
                    })
            except Exception:
                pass
    return issues

def timeline_mismatch_checks(slides_data):
    """
    Collect dates mentioned per slide and find contradictions (e.g., forecasts for 2026 vs chart saying 2025).
    Heuristic: If slide-level dates sets have contradictory ranges or same-named items assigned to different dates.
    """
    slide_dates = {}
    for slide_no, s in slides_data.items():
        text = s["text"] + "\n" + s["ocr_text"]
        dts = extract_dates(text)
        slide_dates[slide_no] = dts

    issues = []

    year_map = defaultdict(list)  
    for slide_no, dts in slide_dates.items():
        for raw, iso in dts:
            try:
                y = int(iso.split('-')[0])
                year_map[y].append((slide_no, raw))
            except:
                pass
    forecast_issues = []
    for slide_no, s in slides_data.items():
        text = (s["text"] + " " + s["ocr_text"]).lower()
        if "forecast" in text or "projection" in text or "expected" in text or "estimated" in text:
            # collect years on this slide
            dts = slide_dates.get(slide_no, [])
            yrs = sorted({int(d[1].split('-')[0]) for d in dts}) if dts else []
            if yrs:
                forecast_issues.append((slide_no, yrs))
    if forecast_issues:
        # simply flag if forecasts refer to disjoint single years across multiple slides
        all_years = set()
        for sl, yrs in forecast_issues:
            all_years.update(yrs)
        if len(all_years) > 1:
            issues.append({
                "message": "Slides contain forecasts/projections mentioning multiple different years",
                "forecast_slide_years": [{"slide": sl, "years": yrs} for sl, yrs in forecast_issues],
                "all_years": sorted(list(all_years))
            })
    return issues

def basic_consistency_checks(slides_data):
    return {
        "number_conflicts": find_conflicting_numbers(slides_data),
        "percentage_issues": percentages_sum_checks(slides_data),
        "timeline_issues": timeline_mismatch_checks(slides_data)
    }



GEMINI_PROMPT_TEMPLATE = """
You are an assistant that finds factual and logical inconsistencies across multiple slides in a presentation.
Input: A list of slides; each slide has: slide_number, text (text extracted from shapes), ocr_text (text read from images), and notes.
Task: For the provided slides, identify potential contradictions or inconsistencies such as:
 - Two slides stating different numeric values for the same metric (e.g., 'Revenue = $10M' vs 'Revenue = $12M')
 - Contradictory textual claims (e.g., 'market is highly competitive' vs 'few competitors')
 - Timeline inconsistencies (dates/forecast years mismatch)
 - Statements that logically contradict each other
Return: A JSON array where each element is an object:
{
  "type": "numeric_conflict" | "text_conflict" | "timeline_conflict" | "uncertain_claim" | "other",
  "slides_involved": [<slide_numbers>],
  "summary": "<short human summary>",
  "evidence": ["<quote from slide or ocr>"],
  "confidence": <0.0-1.0>
}
Only output JSON (no extra commentary).
Be conservative and focus on likely contradictions; include confidence estimates.
"""

def call_gemini_find_contradictions(slides_data, api_key=None, model_name="gemini-2.5-flash", max_chars=60000):
    """
    Sends the slides to Gemini in manageable chunks. Returns model JSON-parsed output or None.
    """
    if gemini_client is None:
        raise RuntimeError("google-genai SDK not installed or importable. pip install google-genai")

    
    slides_list = []
    for slide_no, d in slides_data.items():
        slides_list.append({
            "slide": slide_no,
            "text": (d.get("text","") or "")[:2000],
            "ocr_text": (d.get("ocr_text","") or "")[:2000],
            "notes": (d.get("notes","") or "")[:1000]
        })

    
    client = gemini_client.Client(api_key=api_key) if api_key else gemini_client.Client()
    slides_json = json.dumps(slides_list, ensure_ascii=False)
    if len(slides_json) > max_chars:
        chunks = []
        size = len(slides_json)
        n_chunks = max(1, min(6, (size // max_chars) + 1))
        per = max(1, len(slides_list) // n_chunks)
        for i in range(0, len(slides_list), per):
            chunks.append(slides_list[i:i+per])
    else:
        chunks = [slides_list]

    combined_results = []
    for c in chunks:
        prompt = GEMINI_PROMPT_TEMPLATE + "\n\nSlides:\n" + json.dumps(c, ensure_ascii=False, indent=2)
        # call model
        # using simple generate API
        response = client.generate_text(
            model=model_name,
            temperature=0.0,
            max_output_tokens=1024,
            input=prompt,
        )
        # response.text may have newlines; try parse JSON
        text = response.text.strip()
        # attempt JSON parse; if fails, try to extract JSON substring
        j = None
        try:
            j = json.loads(text)
        except Exception:
            # find first '[' and last ']' and parse substring
            a = text.find('[')
            b = text.rfind(']')
            if a != -1 and b != -1 and b > a:
                try:
                    j = json.loads(text[a:b+1])
                except:
                    j = None
        if j:
            combined_results.extend(j)
        else:
            # if cannot parse, include as 'other' with low confidence
            combined_results.append({
                "type": "other",
                "slides_involved": [],
                "summary": "LLM output unparseable; raw output included in 'evidence'.",
                "evidence": [text],
                "confidence": 0.3
            })
    return combined_results



def analyze_presentation(pptx_path, args):
    slides_data, tmpdir = pptx_to_slides(pptx_path, images_outdir=args.tmpdir, ocr_on_images=(not args.no_ocr))
    report = {
        "file": pptx_path,
        "slides_parsed": len(slides_data),
        "deterministic_checks": basic_consistency_checks(slides_data),
        "llm_checks": None,
        "errors": []
    }

    if not args.no_llm:
        if gemini_client is None:
            report["errors"].append("google-genai SDK not available; skip LLM checks. pip install google-genai")
        else:
            try:
                api_key = args.api_key or os.environ.get("GEMINI_API_KEY")
                llm_res = call_gemini_find_contradictions(slides_data, api_key=api_key, model_name=args.model)
                report["llm_checks"] = llm_res
            except Exception as e:
                report["errors"].append(f"LLM error: {repr(e)}")
    if tmpdir:
        if not args.keep_tmp:
            try:
                shutil.rmtree(tmpdir)
            except:
                pass
        else:
            report["tmpdir"] = tmpdir
    return report

def main():
    parser = argparse.ArgumentParser(description="Find inconsistencies across PPTX slides.")
    parser.add_argument("path_dir", help=".pptx file to analyze")
    parser.add_argument("--no-llm", action="store_true", help="Disable calls to Gemini/LLM (run deterministic checks only)")
    parser.add_argument("gemini_api", default=None, help="Gemini API key (or set GEMINI_API_KEY env var)")
    parser.add_argument("--model", default="gemini-2.5-flash", help="Gemini model name (e.g., gemini-2.5-flash)")
    parser.add_argument("temp_dir", default=None, help="Directory to store extracted images (default: tempdir)")
    parser.add_argument("--keep-tmp", action="store_true", help="Don't delete temp image dir after run")
    parser.add_argument("--no-ocr", action="store_true", help="Don't run OCR on slide images")
    parser.add_argument("--out", default=None, help="Write JSON report to file")
    args = parser.parse_args()

    if not os.path.exists(args.pptx):
        print("ERROR: pptx file not found:", args.pptx)
        return

    try:
        report = analyze_presentation(args.pptx, args)
        text = json.dumps(report, indent=2, ensure_ascii=False)
        print(text)
        if args.out:
            with open(args.out, "w", encoding="utf-8") as f:
                f.write(text)
            print("Report written to", args.out)
    except Exception as e:
        print("Fatal error:", repr(e))

if __name__ == "__main__":
    main()
